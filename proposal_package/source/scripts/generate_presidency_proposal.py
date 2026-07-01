from datetime import date
from math import ceil
from pathlib import Path
from textwrap import wrap

from PIL import Image as PILImage, ImageChops

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import (
    BaseDocTemplate,
    Frame,
    PageBreak,
    PageTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
    Image,
    KeepTogether,
)

OUTPUT_DIR = Path(__file__).resolve().parent
TODAY = date.today().isoformat()
DOC_ID = "BT-PRESIDENCY-PROPOSAL"
COMPANY = "Bluetrace Technologies Pvt. Ltd."
SUBTITLE = "Technology Innovation Partner for Presidency English School, Parbhani"
PDF_NAME = "Bluetrace_Presidency_Principal_Director_Management_Proposal_v1.pdf"
PPTX_NAME = "Bluetrace_Presidency_Principal_Director_Management_Deck_v1.pptx"
LOGO_CANDIDATES = ["bluetrace_logo.png", "logo.png", "bluetrace-logo.png"]
STARTUP_LOGO_CANDIDATES = ["startup_india_logo.png", "dpiit_logo.png", "startupindia_logo.png"]
UDYAM_LOGO_CANDIDATES = ["udyam_logo.png", "msme_logo.png"]
MCA_LOGO_CANDIDATES = ["mca_logo.png", "corporate_affairs_logo.png"]
EXTERNAL_ASSET_DIRS = [
    Path(r"C:\Users\ansar\Downloads\blutrace"),
]
CREDENTIAL_REFERENCES = [
    ("Certificate of Incorporation", "INC-34_1-25411283954.pdf"),
    ("DPIIT Recognition", "DIPP266429_BLUETRACE_TECHNOLOGIES_PRIVATE_LIMITED_RECOGNITION_2763114265465494817.pdf"),
    ("Udyam Registration", "Print _ Udyam Registration Certificate.pdf"),
    ("Udyam Annexure", "Print _ Udyam Registration Certificate ANNX.pdf"),
    ("SPICE+ Part B Approval", "SPICE + Part B_Approval Letter_AC3221601.pdf"),
    ("Company Letterhead", "IT and Software Letterhead #19223.pdf"),
]

NAVY = colors.HexColor("#14324a")
TEAL = colors.HexColor("#1f8a70")
GOLD = colors.HexColor("#c58f2b")
LIGHT = colors.HexColor("#eef4f7")
SOFT = colors.HexColor("#f7fafc")
DARK = colors.HexColor("#24333f")
MUTED = colors.HexColor("#5c6b73")
BORDER = colors.HexColor("#d5dee5")
ACCENT = colors.HexColor("#e9f6f2")
PPTX_NAVY = RGBColor(20, 50, 74)
PPTX_TEAL = RGBColor(31, 138, 112)
PPTX_GOLD = RGBColor(197, 143, 43)
PPTX_WHITE = RGBColor(255, 255, 255)


def find_logo_path() -> Path | None:
    search_dirs = [OUTPUT_DIR] + EXTERNAL_ASSET_DIRS
    for base in search_dirs:
        for name in LOGO_CANDIDATES:
            candidate = base / name
            if candidate.exists():
                return candidate
    return None


def _find_first_existing(candidates: list[str]) -> Path | None:
    search_dirs = [OUTPUT_DIR] + EXTERNAL_ASSET_DIRS
    for base in search_dirs:
        for name in candidates:
            candidate = base / name
            if candidate.exists():
                return candidate
    return None


def find_credential_logo_paths() -> list[Path]:
    logos: list[Path] = []
    startup = _find_first_existing(STARTUP_LOGO_CANDIDATES)
    udyam = _find_first_existing(UDYAM_LOGO_CANDIDATES)
    mca = _find_first_existing(MCA_LOGO_CANDIDATES)
    for p in [startup, udyam, mca]:
        if p is not None:
            logos.append(get_trimmed_logo_path(p))
    return logos


def get_trimmed_logo_path(path: Path) -> Path:
    """Trim white/transparent borders around a logo for cleaner placement."""
    cache_dir = OUTPUT_DIR / ".logo_cache"
    cache_dir.mkdir(exist_ok=True)
    trimmed_path = cache_dir / f"{path.stem}_trimmed.png"
    if trimmed_path.exists() and trimmed_path.stat().st_mtime >= path.stat().st_mtime:
        return trimmed_path

    try:
        with PILImage.open(path) as img:
            rgba = img.convert("RGBA")
            bg = PILImage.new("RGBA", rgba.size, (255, 255, 255, 255))
            diff = ImageChops.difference(rgba, bg)
            bbox = diff.getbbox()
            if not bbox:
                rgba.save(trimmed_path, "PNG")
                return trimmed_path

            pad = 6
            left = max(0, bbox[0] - pad)
            top = max(0, bbox[1] - pad)
            right = min(rgba.width, bbox[2] + pad)
            bottom = min(rgba.height, bbox[3] + pad)
            cropped = rgba.crop((left, top, right, bottom))
            cropped.save(trimmed_path, "PNG")
            return trimmed_path
    except Exception:
        return path


def fit_in_box(img_w: float, img_h: float, box_w: float, box_h: float) -> tuple[float, float]:
    """Return width/height that preserves aspect ratio within a target box."""
    if img_w <= 0 or img_h <= 0:
        return box_w, box_h
    scale = min(box_w / img_w, box_h / img_h)
    return img_w * scale, img_h * scale


def place_logos_bottom_aligned(slide, logos: list[Path], left: float, top: float, width: float, height: float, gap: float) -> None:
    """Place logos in a horizontal strip with shared bottom baseline and centered group."""
    if not logos:
        return

    max_logo_h = height * 0.76
    sizes: list[tuple[Path, float, float]] = []
    for logo in logos:
        try:
            with PILImage.open(logo) as logo_img:
                w, h = fit_in_box(float(logo_img.width), float(logo_img.height), width, max_logo_h)
            sizes.append((logo, w, h))
        except Exception:
            continue

    if not sizes:
        return

    total_width = sum(w for _, w, _ in sizes) + (len(sizes) - 1) * gap
    x = left + max(0.0, (width - total_width) / 2.0)
    baseline = top + height - 0.03

    for logo, w, h in sizes:
        y = baseline - h
        slide.shapes.add_picture(str(logo), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        x += w + gap


def estimate_content_box_height(bullets: list[str], has_logo_strip: bool) -> float:
    """Estimate a balanced content-box height in inches based on bullet density."""
    line_count = 0
    for bullet in bullets:
        line_count += max(1, ceil(len(bullet) / 54))

    # Base + per-line growth tuned for 16:9 slide readability.
    estimated = 1.45 + (line_count * 0.44)
    if has_logo_strip:
        estimated += 0.3
    return max(2.8, min(4.9, estimated))


def credential_badge_labels() -> list[str]:
    return [
        "Startup India (DPIIT)",
        "MSME/Udyam",
        "MCA Registered",
    ]


def available_credential_references() -> list[str]:
    available: list[str] = []
    for label, filename in CREDENTIAL_REFERENCES:
        if (OUTPUT_DIR / filename).exists():
            available.append(f"{label}: verified and available")
            continue
        for base in EXTERNAL_ASSET_DIRS:
            if (base / filename).exists():
                available.append(f"{label}: verified and available")
                break
    return available


def build_pdf() -> None:
    pdf_path = OUTPUT_DIR / PDF_NAME
    doc = BaseDocTemplate(
        str(pdf_path),
        pagesize=A4,
        leftMargin=18 * mm,
        rightMargin=18 * mm,
        topMargin=24 * mm,
        bottomMargin=20 * mm,
        title=SUBTITLE,
        author=COMPANY,
        subject="Presidency English School proposal",
    )

    frame = Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height, id="normal")
    doc.addPageTemplates([
        PageTemplate(id="proposal", frames=[frame], onPage=draw_pdf_header_footer),
    ])

    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        name="Hero",
        parent=styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=22,
        leading=26,
        textColor=NAVY,
        alignment=TA_CENTER,
        spaceAfter=10,
    ))
    styles.add(ParagraphStyle(
        name="HeroSub",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=11,
        leading=15,
        textColor=MUTED,
        alignment=TA_CENTER,
        spaceAfter=12,
    ))
    styles.add(ParagraphStyle(
        name="SectionTitle",
        parent=styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=15,
        leading=18,
        textColor=NAVY,
        spaceBefore=8,
        spaceAfter=8,
    ))
    styles.add(ParagraphStyle(
        name="BodyCopy",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=10,
        leading=14,
        textColor=DARK,
        spaceAfter=6,
    ))
    styles.add(ParagraphStyle(
        name="SmallCopy",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=8.5,
        leading=11,
        textColor=MUTED,
    ))
    styles.add(ParagraphStyle(
        name="BulletCopy",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=10,
        leading=13,
        textColor=DARK,
        leftIndent=12,
        bulletIndent=0,
        spaceAfter=3,
    ))

    story = []
    story.append(Spacer(1, 24))
    story.append(Paragraph(COMPANY, styles["Hero"]))
    story.append(Paragraph(SUBTITLE, styles["HeroSub"]))
    cred_logo_paths = find_credential_logo_paths()
    if cred_logo_paths:
        logo_cells = []
        for p in cred_logo_paths:
            try:
                with PILImage.open(p) as logo_img:
                    w, h = fit_in_box(float(logo_img.width), float(logo_img.height), 28 * mm, 10 * mm)
                logo_cells.append(Image(str(p), width=w, height=h))
            except Exception:
                pass
        if logo_cells:
            logos_table = Table([logo_cells], colWidths=[32 * mm] * len(logo_cells))
            logos_table.setStyle(TableStyle([
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("VALIGN", (0, 0), (-1, -1), "BOTTOM"),
                ("LEFTPADDING", (0, 0), (-1, -1), 3),
                ("RIGHTPADDING", (0, 0), (-1, -1), 3),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 1),
            ]))
            story.append(logos_table)
            story.append(Spacer(1, 6))
    else:
        labels = credential_badge_labels()
        badge_cells = [Paragraph(f"<b>{label}</b>", styles["SmallCopy"]) for label in labels]
        badge_table = Table([badge_cells], colWidths=[44 * mm, 38 * mm, 38 * mm])
        badge_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), LIGHT),
            ("BOX", (0, 0), (-1, -1), 0.6, BORDER),
            ("INNERGRID", (0, 0), (-1, -1), 0.4, BORDER),
            ("ALIGN", (0, 0), (-1, -1), "CENTER"),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ]))
        story.append(badge_table)
        story.append(Spacer(1, 6))
    story.append(Spacer(1, 10))

    cover_table = Table([
        [
            Paragraph("Document ID<br/><b>{}</b>".format(DOC_ID), styles["BodyCopy"]),
            Paragraph("Version<br/><b>v1.0</b>", styles["BodyCopy"]),
            Paragraph("Revision<br/><b>R0</b>", styles["BodyCopy"]),
            Paragraph(f"Date<br/><b>{TODAY}</b>", styles["BodyCopy"]),
        ]
    ], colWidths=[doc.width / 4.0] * 4)
    cover_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), LIGHT),
        ("BOX", (0, 0), (-1, -1), 0.6, BORDER),
        ("INNERGRID", (0, 0), (-1, -1), 0.4, BORDER),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("TOPPADDING", (0, 0), (-1, -1), 10),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
    ]))
    story.append(cover_table)
    story.append(Spacer(1, 12))
    story.append(Paragraph(
        "This proposal positions Bluetrace as a long-term innovation partner for Presidency English School, Parbhani. The strategy focuses on AI, robotics, IoT, smart campus solutions, and teacher productivity rather than commodity website, ERP, or generic training services.",
        styles["BodyCopy"],
    ))
    story.append(Spacer(1, 18))

    credential_bullets = [
        "Bluetrace Technologies Pvt. Ltd. is an incorporated private limited company.",
        "DPIIT Startup India recognition and MSME/Udyam credentials are available for institutional validation.",
        "Official approvals, registration certificates, and company letterhead are maintained in compliance records.",
        "Detailed supporting documents can be shared during due diligence on request.",
    ]
    credential_refs = available_credential_references()
    if credential_refs:
        credential_bullets.extend(credential_refs)

    add_section(story, styles, "0. Corporate Credentials and Compliance", credential_bullets)

    add_section(story, styles, "1. Strategic Positioning", [
        "Bluetrace should lead with transformation, not software delivery.",
        "The offer should feel board-level, future-facing, and student-outcome oriented.",
        "Start with a zero-cost workshop to create trust and management confidence.",
        "Move the conversation toward a recurring annual partnership model.",
    ])

    add_section(story, styles, "2. Current Challenges in Schools", [
        "Teachers: time-heavy question paper preparation, repetitive worksheets, and manual lesson planning.",
        "Students: passive learning patterns, limited practical AI/robotics exposure, and low project-based innovation.",
        "Parents: fragmented updates across WhatsApp, calls, circulars, and social media channels.",
        "Management: no integrated analytics layer, high reporting overhead, and slower data-driven decisions.",
        "Conclusion: Bluetrace addresses these issues through a unified Presidency Digital School Platform.",
    ])

    add_section(story, styles, "3. Presidency Digital School Platform", [
        "The solution is not only a mobile app; it is the school's secure digital backbone.",
        "Core flow: Parent/Student/Teacher apps -> School Cloud Server -> Central Database -> Principal and Management Dashboard.",
        "Platform outcomes: attendance, homework, fees, circulars, admissions, bus tracking, and AI-enabled workflows in one system.",
        "This reduces communication confusion and creates one trusted source of truth for all stakeholders.",
    ])

    add_section(story, styles, "4. How the Platform Works (Operational Flows)", [
        "Homework flow: Teacher uploads homework -> notification engine -> parent/student app -> submission -> teacher review -> dashboard update.",
        "Attendance flow: Teacher marks attendance -> central database sync -> parent notification -> management analytics.",
        "Fee flow: fee due generation -> parent reminder -> payment confirmation -> admin ledger and dashboard update.",
        "Circular flow: office publishes circular -> app push -> SMS/WhatsApp fallback -> parent acknowledgement.",
    ])

    add_section(story, styles, "5. Role-Based Digital Experience", [
        "Parent app: attendance, homework, circulars, fee status, exam timetable, results, bus location, leave requests, teacher messages.",
        "Teacher app: attendance, homework upload, notices, marks entry, AI worksheet generation, lesson planning, parent communication.",
        "Student app: homework, learning videos, quizzes, assignment submission, certificates, AI/robotics activities, project portfolio.",
        "Principal and management dashboard: attendance trends, fee collection, admissions, teacher activity, engagement, transport, and complaints.",
    ])

    add_section(story, styles, "5A. Smart Wearable Student Safety Layer (Optional Pilot)", [
        "Use case: optional Bluetooth wearable tag for transport and campus transitions to improve child visibility and parent trust.",
        "School bus flow: boarding and drop events can trigger parent app confirmation updates.",
        "Campus flow: classroom, library, and pickup zones can be mapped as event points without invasive continuous tracking.",
        "Safety by design: event-based monitoring, consent workflow, role-based visibility, and policy-aligned retention windows.",
        "Business value: stronger safety positioning, lower parent anxiety, and premium institutional differentiation.",
    ])

    add_section(story, styles, "6. Admin Control and Notification Backbone", [
        "Admin panel controls students, parents, teachers, fees, timetable, admissions, transport, circulars, and reports.",
        "Notification strategy: app push first, then SMS/WhatsApp backup when needed.",
        "Every major communication can be acknowledgement-tracked to improve follow-through.",
        "This removes dependency on scattered message threads and missed updates.",
    ])

    add_section(story, styles, "7. Engineering Depth and Technical Credibility", [
        "Bluetrace engineering stack includes Android, embedded Linux, STM32, nRF52, ESP32, BLE, PCB design, and cloud integrations.",
        "Capability extends from app workflows to device integrations such as RFID, bus GPS, and smart campus sensors.",
        "The same engineering team can deliver AI pipelines, IoT gateways, and production-grade connected systems.",
        "This is a major differentiator versus training-only or software-only vendors.",
    ])

    add_section(story, styles, "8. Presidency Vision 2030 Roadmap", [
        "2026: Teacher AI productivity rollout and unified communication platform.",
        "2027: Student AI learning modules, robotics foundation, and project showcases.",
        "2028: Robotics and IoT innovation lab with curriculum-linked practical programs.",
        "2029: Smart campus expansion (attendance, transport, energy, and integrated dashboards).",
        "2030: AI analytics-driven school operations and innovation-centered institutional positioning.",
    ])

    add_section(story, styles, "9. Investment Logic: Why Management Should Invest", [
        "AI teacher tools -> saves planning and assessment effort, increasing classroom quality time.",
        "Digital school platform -> faster parent communication and reduced repetitive office workload.",
        "Robotics/AI programs -> stronger student outcomes and practical STEM profile.",
        "Smart operations dashboard -> better management visibility and faster decisions.",
        "Institutional outcome -> stronger school reputation and clearer admissions differentiation.",
    ])

    revenue_heading = Paragraph("10. Revenue / Partnership Model", styles["SectionTitle"])
    revenue_rows = [
        ["Option", "Scope", "Indicative Value"],
        ["A", "Teacher training", "₹20,000 - ₹50,000"],
        ["B", "AI program for students", "₹500 per student per month"],
        ["C", "Innovation lab setup", "₹2 lakh - ₹5 lakh"],
        ["D", "Annual technology partnership", "₹3 lakh - ₹10 lakh per year"],
    ]
    revenue_table = Table(revenue_rows, colWidths=[20 * mm, 95 * mm, 50 * mm])
    revenue_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), NAVY),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("BACKGROUND", (0, 1), (-1, -1), SOFT),
        ("BOX", (0, 0), (-1, -1), 0.6, BORDER),
        ("INNERGRID", (0, 0), (-1, -1), 0.4, BORDER),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("LEADING", (0, 0), (-1, -1), 11),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]))
    story.append(KeepTogether([revenue_heading, revenue_table, Spacer(1, 10)]))

    add_section(story, styles, "11. 90-Day Pilot Roadmap", [
        "Weeks 1-2: free teacher workshop and management alignment.",
        "Weeks 3-6: platform pilot workflows for homework, attendance, and parent updates.",
        "Weeks 7-10: AI student module and innovation club launch with tracked outcomes; optional wearable safety pilot in one route/zone.",
        "Weeks 11-12: management review, KPI snapshot, and annual partnership proposal.",
    ])

    add_section(story, styles, "12. Student Storyboard (Cartoon Narrative for Management Pitch)", [
        "Panel 1 - Morning rush: parent uncertainty before bus pickup creates stress and repeated calls.",
        "Panel 2 - Smart event: wearable + app confirms safe boarding in real time.",
        "Panel 3 - School day: homework and circular updates arrive in one trusted app channel.",
        "Panel 4 - Issue moment: delayed route or class change triggers guided alert and acknowledgement.",
        "Panel 5 - End of day: pickup confirmation closes loop for parent, school, and management dashboard.",
        "Panel 6 - Outcome view: fewer missed communications, faster escalation, and better parent confidence.",
    ])

    add_section(story, styles, "13. Why Bluetrace vs Generic Vendors", [
        "Generic vendors mostly provide standard training; Bluetrace provides engineering plus training.",
        "Generic vendors are curriculum-limited; Bluetrace designs customized workflows, labs, and technical integrations.",
        "Generic vendors are often one-time service providers; Bluetrace is positioned as a long-term technology partner.",
        "Bluetrace combines app, AI, embedded, IoT, and cloud capabilities under one delivery model.",
    ])

    story.append(KeepTogether([
        Paragraph("14. Next Step", styles["SectionTitle"]),
        Paragraph(
            "Use this as the board-facing proposal that answers why to invest, what to deploy first, and how value will be measured. The recommended follow-up is a chairman/principal workshop and a 90-day pilot signoff focused on measurable adoption outcomes.",
            styles["BodyCopy"],
        ),
        Spacer(1, 8),
    ]))
    story.append(Paragraph(
        "Confidentiality: Internal business development draft prepared for client discussion purposes.",
        styles["SmallCopy"],
    ))

    doc.build(story)


def add_section(story, styles, title, bullets):
    section_title = Paragraph(title, styles["SectionTitle"])
    data = []
    for bullet in bullets:
        data.append([Paragraph(f"• {bullet}", styles["BodyCopy"])])
    table = Table(data, colWidths=[168 * mm])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), ACCENT),
        ("BOX", (0, 0), (-1, -1), 0.6, BORDER),
        ("INNERGRID", (0, 0), (-1, -1), 0.35, BORDER),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 8),
        ("RIGHTPADDING", (0, 0), (-1, -1), 8),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]))
    story.append(KeepTogether([section_title, table, Spacer(1, 8)]))


def draw_pdf_header_footer(canvas, doc):
    canvas.saveState()
    width, height = A4
    canvas.setFillColor(NAVY)
    canvas.rect(0, height - 18 * mm, width, 18 * mm, fill=1, stroke=0)

    logo_path = find_logo_path()
    if logo_path is not None:
        try:
            canvas.drawImage(
                str(logo_path),
                doc.leftMargin,
                height - 16.5 * mm,
                width=14 * mm,
                height=14 * mm,
                preserveAspectRatio=True,
                mask="auto",
            )
            title_x = doc.leftMargin + 16 * mm
        except Exception:
            title_x = doc.leftMargin
    else:
        title_x = doc.leftMargin

    canvas.setFillColor(colors.white)
    canvas.setFont("Helvetica-Bold", 11)
    canvas.drawString(title_x, height - 11.5 * mm, COMPANY)
    canvas.setFont("Helvetica", 8.5)
    canvas.drawRightString(width - doc.rightMargin, height - 11.5 * mm, DOC_ID)
    canvas.setFillColor(GOLD)
    canvas.setFont("Helvetica-Bold", 9)
    canvas.drawString(title_x, height - 15.8 * mm, SUBTITLE)

    canvas.setFillColor(NAVY)
    canvas.rect(0, 0, width, 13 * mm, fill=1, stroke=0)
    canvas.setFillColor(colors.white)
    canvas.setFont("Helvetica", 8)
    canvas.drawString(doc.leftMargin, 5 * mm, f"Confidential - Internal Proposal Draft")
    canvas.drawRightString(width - doc.rightMargin, 5 * mm, f"Page {canvas.getPageNumber()}")
    canvas.restoreState()


def build_pptx() -> None:
    pptx_path = OUTPUT_DIR / PPTX_NAME
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = [
        {
            "title": "Strategic Proposal",
            "subtitle": "Technology Innovation Partner for Presidency English School, Parbhani",
            "bullets": [
                "Board-level proposal for AI, robotics, IoT, and smart campus transformation",
                "Positioned for long-term partnership rather than commodity software delivery",
            ],
            "footer": "Presidency English School Proposal | Confidential",
            "kind": "hero",
        },
        {
            "title": "Why This Approach Works",
            "bullets": [
                "The school already has digital boards, internet, online admissions, and online fee payment.",
                "That means the market message should move beyond generic IT services.",
                "The winning story is future-skills, innovation, and measurable student value.",
            ],
            "footer": "Strategic Positioning",
            "kind": "content",
        },
        {
            "title": "Corporate Credentials and Compliance",
            "bullets": [
                "Bluetrace Technologies Pvt. Ltd. is an incorporated private limited company.",
                "Startup India (DPIIT) recognition and MSME/Udyam credentials are in place.",
                "Incorporation and statutory approvals are maintained for institutional due diligence.",
                "Official company letterhead and supporting compliance records are available on request.",
                "Proposal deck intentionally uses clean credential summaries instead of raw file names.",
            ],
            "footer": "Trust and Compliance",
            "kind": "content",
        },
        {
            "title": "Priority Focus Areas",
            "bullets": [
                "Answer management's core question first: why should the school invest now?",
                "Frame value in admissions, parent satisfaction, teacher productivity, and reputation.",
                "Position Bluetrace as a long-term transformation partner, not a software vendor.",
                "Use pilot-to-partnership progression with measurable outcomes.",
            ],
            "footer": "Priority strategy",
            "kind": "content",
        },
        {
            "title": "Current Challenges in Schools",
            "bullets": [
                "Teachers: repetitive planning, worksheets, and paper generation load",
                "Students: limited AI/robotics exposure and low project-based innovation",
                "Parents: scattered communication channels and missed updates",
                "Management: limited analytics and heavy manual reporting",
            ],
            "footer": "Problem statement",
            "kind": "content",
        },
        {
            "title": "Presidency Digital School Platform",
            "bullets": [
                "Present as a platform, not only a mobile app",
                "Parent/Student/Teacher apps -> cloud server -> central database",
                "Admin + Principal dashboard for decisions, control, and accountability",
                "One secure digital backbone across communication, academics, and operations",
            ],
            "footer": "Platform architecture",
            "kind": "content",
        },
        {
            "title": "How It Works: End-to-End Operational Flow",
            "bullets": [
                "Homework: teacher upload -> cloud -> notification -> parent/student -> submission -> review",
                "Attendance: teacher mark -> database sync -> parent alert -> dashboard analytics",
                "Fees: due generation -> reminder -> payment confirmation -> admin ledger update",
                "Circulars/notices: publish -> app push -> SMS/WhatsApp fallback -> acknowledgement tracking",
            ],
            "footer": "Workflow confidence",
            "kind": "content",
        },
        {
            "title": "Role-Based User Experience",
            "bullets": [
                "Parent app: attendance, homework, fees, results, circulars, bus tracking, leave",
                "Teacher app: attendance, homework, marks, notices, parent communication, AI worksheet",
                "Student app: assignments, quizzes, videos, portfolio, certificates, innovation tasks",
                "Admin panel: students, teachers, timetable, transport, reports, admissions",
            ],
            "footer": "Experience design",
            "kind": "content",
        },
        {
            "title": "Optional Differentiator: Smart Wearable Safety Layer",
            "bullets": [
                "Event-based wearable model for bus boarding, transit, and pickup checkpoints",
                "Parent app receives trusted safety updates without invasive continuous surveillance",
                "Principal dashboard tracks route-level exceptions and response times",
                "Pilot can start with one grade, one route, and policy-aligned consent",
            ],
            "footer": "Wearable safety use case",
            "kind": "content",
        },
        {
            "title": "Management Dashboard and Decision Layer",
            "bullets": [
                "Daily attendance trends and class-wise performance visibility",
                "Fee collection status and pending follow-up indicators",
                "Admission enquiry pipeline and parent engagement insights",
                "Teacher activity, bus status, issues, and escalation monitoring",
            ],
            "footer": "Executive control",
            "kind": "content",
        },
        {
            "title": "Notification Backbone (No More Channel Confusion)",
            "bullets": [
                "School publishes notice once from central panel",
                "App notification reaches right audience instantly",
                "SMS/WhatsApp fallback triggers only when required",
                "Parent acknowledgement becomes measurable and auditable",
            ],
            "footer": "Reliable communication",
            "kind": "content",
        },
        {
            "title": "Student Storyboard for Parent Trust",
            "bullets": [
                "Story panel 1: parent uncertainty before morning transport",
                "Story panel 2: smart wearable/app confirms safe boarding",
                "Story panel 3: homework and notices delivered in one app",
                "Story panel 4: delay alert and acknowledgement close communication gaps",
                "Story panel 5: pickup confirmation and confidence dashboard outcome",
            ],
            "footer": "Narrative communication",
            "kind": "content",
        },
        {
            "title": "Presidency Vision 2030",
            "bullets": [
                "2026: Teacher AI productivity + unified communication baseline",
                "2027: Student AI modules and project-based learning scale-up",
                "2028: Robotics/IoT lab with curriculum-linked practical outcomes",
                "2029-2030: Smart campus + AI analytics + innovation center positioning",
            ],
            "footer": "Long-term roadmap",
            "kind": "content",
        },
        {
            "title": "Revenue / Partnership Models",
            "bullets": [
                "Teacher training: ₹20,000 - ₹50,000",
                "AI program: ₹500 per student per month",
                "Innovation lab: ₹2 lakh - ₹5 lakh setup",
                "Annual technology partnership: ₹3 lakh - ₹10 lakh per year",
            ],
            "footer": "Commercial options",
            "kind": "content",
        },
        {
            "title": "ROI and Outcome Mapping",
            "bullets": [
                "AI teacher tools -> reduced planning effort and better instructional quality time",
                "Digital platform -> fewer repetitive office calls and faster parent communication",
                "Robotics/AI activities -> stronger STEM profile and practical student outcomes",
                "Dashboard analytics -> better management decisions with integrated school data",
            ],
            "footer": "Investment justification",
            "kind": "content",
        },
        {
            "title": "90-Day Pilot Implementation Path",
            "bullets": [
                "Weeks 1-2: leadership workshop, baseline metrics, and pilot scope signoff",
                "Weeks 3-6: core platform workflows (homework, attendance, communication, fees)",
                "Weeks 7-10: AI modules plus optional wearable safety pilot (single route)",
                "Weeks 11-12: KPI review, board presentation, and annual partnership decision",
            ],
            "footer": "Execution roadmap",
            "kind": "content",
        },
        {
            "title": "Why Bluetrace vs Generic Vendor",
            "bullets": [
                "Generic vendor: training-only model | Bluetrace: engineering + training model",
                "Generic vendor: standard curriculum | Bluetrace: customized platform + workflows",
                "Generic vendor: limited hardware depth | Bluetrace: STM32, nRF52, ESP32, BLE, PCB, IoT",
                "Generic vendor: one-time engagement | Bluetrace: long-term technology partnership",
            ],
            "footer": "Differentiation",
            "kind": "content",
        },
        {
            "title": "Next Step",
            "bullets": [
                "Use this deck for a chairman/principal meeting",
                "Lead with the free workshop and pilot offer, then storyboard walkthrough",
                "Close with an annual partnership model, not a one-off project",
            ],
            "footer": "Final call to action",
            "kind": "closing",
        },
    ]

    for index, slide_info in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        add_ppt_background(slide)
        add_ppt_header(slide, slide_info["title"], index, len(slides))
        if slide_info["kind"] == "hero":
            add_ppt_hero(slide, slide_info)
        else:
            add_ppt_content(slide, slide_info)
        add_ppt_footer(slide, slide_info["footer"])

    if prs.slides and len(prs.slides) > len(slides):
        xml_slides = prs.slides._sldIdLst  # type: ignore[attr-defined]
        first = xml_slides[0]
        xml_slides.remove(first)

    prs.core_properties.title = SUBTITLE
    prs.core_properties.subject = "Presidency English School proposal"
    prs.core_properties.author = COMPANY
    prs.core_properties.comments = "Generated by Bluetrace proposal generator"
    prs.save(str(pptx_path))


def add_ppt_background(slide):
    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    left.fill.solid()
    left.fill.fore_color.rgb = PPTX_NAVY
    left.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.8), 0, Inches(2.533), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PPTX_TEAL
    accent.line.fill.background()

    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.8), Inches(-0.2), Inches(2.4), Inches(2.4))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = PPTX_GOLD
    circle1.line.fill.background()
    circle1.fill.transparency = 0.35

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.6), Inches(5.8), Inches(2.2), Inches(2.2))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = PPTX_GOLD
    circle2.line.fill.background()
    circle2.fill.transparency = 0.75


def add_ppt_header(slide, title, slide_no, total):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(0.25), Inches(12.6), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PPTX_WHITE
    bar.line.fill.background()
    bar.fill.transparency = 0.06

    logo_path = find_logo_path()
    if logo_path is not None:
        try:
            slide.shapes.add_picture(str(logo_path), Inches(0.42), Inches(0.28), width=Inches(0.34), height=Inches(0.34))
            company_x = Inches(0.82)
            company_w = Inches(8.18)
        except Exception:
            company_x = Inches(0.5)
            company_w = Inches(8.5)
    else:
        company_x = Inches(0.5)
        company_w = Inches(8.5)

    tx = slide.shapes.add_textbox(company_x, Inches(0.3), company_w, Inches(0.4))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = COMPANY
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = PPTX_NAVY

    tx2 = slide.shapes.add_textbox(Inches(9.2), Inches(0.3), Inches(3.1), Inches(0.4))
    tf2 = tx2.text_frame
    tf2.text = f"{DOC_ID}  |  {slide_no}/{total}"
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tf2.paragraphs[0].runs[0].font.size = Pt(9)
    tf2.paragraphs[0].runs[0].font.color.rgb = PPTX_NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = PPTX_GOLD
    line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(9.5), Inches(0.8))
    ttf = title_box.text_frame
    ttf.word_wrap = True
    ttf.vertical_anchor = MSO_ANCHOR.TOP
    p = ttf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = PPTX_WHITE


def add_ppt_footer(slide, footer_text):
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(7.0), Inches(12.6), Inches(0.32))
    footer.fill.solid()
    footer.fill.fore_color.rgb = PPTX_WHITE
    footer.line.fill.background()
    footer.fill.transparency = 0.08

    left = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(5.8), Inches(0.2))
    tf = left.text_frame
    tf.text = footer_text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.runs[0].font.size = Pt(8)
    p.runs[0].font.color.rgb = PPTX_NAVY

    right = slide.shapes.add_textbox(Inches(8.7), Inches(7.05), Inches(4.0), Inches(0.2))
    tf2 = right.text_frame
    tf2.text = f"Confidential | {TODAY}"
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    p2.runs[0].font.size = Pt(8)
    p2.runs[0].font.color.rgb = PPTX_NAVY


def add_ppt_hero(slide, slide_info):
    content_left = 0.9
    content_width = 6.9

    subtitle_box = slide.shapes.add_textbox(Inches(content_left), Inches(2.85), Inches(content_width), Inches(0.9))
    tf2 = subtitle_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = slide_info["subtitle"]
    r2.font.size = Pt(15)
    r2.font.color.rgb = PPTX_WHITE

    logos = find_credential_logo_paths()
    strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(content_left), Inches(3.8), Inches(content_width), Inches(0.52))
    strip.fill.solid()
    strip.fill.fore_color.rgb = PPTX_WHITE
    strip.line.color.rgb = PPTX_WHITE
    strip.fill.transparency = 0.08

    if logos:
        place_logos_bottom_aligned(slide, logos, left=1.08, top=3.86, width=6.45, height=0.38, gap=0.28)
    else:
        x = 1.08
        for label in credential_badge_labels():
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.86), Inches(1.95), Inches(0.3))
            badge.fill.solid()
            badge.fill.fore_color.rgb = PPTX_WHITE
            badge.line.color.rgb = PPTX_NAVY
            txt = slide.shapes.add_textbox(Inches(x + 0.02), Inches(3.91), Inches(1.9), Inches(0.18))
            tfb = txt.text_frame
            tfb.text = label
            p = tfb.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.size = Pt(8)
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = PPTX_NAVY
            x += 2.15

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(content_left), Inches(4.4), Inches(content_width), Inches(1.05))
    box.fill.solid()
    box.fill.fore_color.rgb = PPTX_WHITE
    box.fill.transparency = 0.08
    box.line.color.rgb = PPTX_WHITE

    tx = slide.shapes.add_textbox(Inches(1.1), Inches(4.52), Inches(6.45), Inches(0.8))
    tf3 = tx.text_frame
    tf3.word_wrap = True
    lines = [
        "AI for teachers",
        "AI excellence for students",
        "Robotics and IoT innovation lab",
        "Smart campus solutions",
    ]
    first = True
    for line in lines:
        p = tf3.paragraphs[0] if first else tf3.add_paragraph()
        first = False
        p.text = f"• {line}"
        p.font.size = Pt(12)
        p.font.color.rgb = PPTX_NAVY

    right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.55), Inches(1.75), Inches(3.7), Inches(3.9))
    right.fill.solid()
    right.fill.fore_color.rgb = PPTX_WHITE
    right.line.color.rgb = PPTX_WHITE
    right.fill.transparency = 0.08

    txr = slide.shapes.add_textbox(Inches(8.8), Inches(2.0), Inches(3.1), Inches(3.3))
    tfr = txr.text_frame
    tfr.word_wrap = True
    intro = [
        "Presidency English School already has digital foundations.",
        "Bluetrace should now position the next layer of differentiation.",
        "This deck is designed for trustees, management, and principal-level discussion.",
    ]
    first = True
    for item in intro:
        p = tfr.paragraphs[0] if first else tfr.add_paragraph()
        first = False
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = PPTX_NAVY
        p.space_after = Pt(10)


def add_ppt_content(slide, slide_info):
    has_logo_strip = slide_info.get("title", "").strip().lower() == "corporate credentials and compliance"
    box_h = estimate_content_box_height(slide_info["bullets"], has_logo_strip)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.55), Inches(9.0), Inches(box_h))
    box.fill.solid()
    box.fill.fore_color.rgb = PPTX_WHITE
    box.line.color.rgb = PPTX_WHITE
    box.fill.transparency = 0.04

    content_y = Inches(1.85)
    content_h = Inches(box_h - 0.65)

    if has_logo_strip:
        logos = find_credential_logo_paths()
        strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(1.86), Inches(8.4), Inches(0.58))
        strip.fill.solid()
        strip.fill.fore_color.rgb = PPTX_WHITE
        strip.line.color.rgb = PPTX_WHITE
        strip.fill.transparency = 0.0

        if logos:
            place_logos_bottom_aligned(slide, logos, left=1.22, top=1.93, width=7.9, height=0.43, gap=0.22)
        else:
            x = 1.22
            for label in credential_badge_labels():
                badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.95), Inches(2.1), Inches(0.35))
                badge.fill.solid()
                badge.fill.fore_color.rgb = PPTX_WHITE
                badge.line.color.rgb = PPTX_NAVY
                txt = slide.shapes.add_textbox(Inches(x + 0.03), Inches(2.01), Inches(2.04), Inches(0.2))
                tfb = txt.text_frame
                tfb.text = label
                p = tfb.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.runs[0].font.size = Pt(9)
                p.runs[0].font.bold = True
                p.runs[0].font.color.rgb = PPTX_NAVY
                x += 2.25

        content_y = Inches(2.55)
        content_h = Inches(max(1.8, box_h - 1.05))

    tx = slide.shapes.add_textbox(Inches(1.05), content_y, Inches(8.4), content_h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for bullet in slide_info["bullets"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18 if len(bullet) < 48 else 14)
        p.font.color.rgb = PPTX_NAVY
        p.space_after = Pt(8)

    sidebar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.15), Inches(1.55), Inches(2.35), Inches(box_h))
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = PPTX_WHITE
    sidebar.line.color.rgb = PPTX_WHITE
    sidebar.fill.transparency = 0.05

    side_text = slide.shapes.add_textbox(Inches(10.35), Inches(1.85), Inches(1.95), Inches(max(1.8, box_h - 0.8)))
    stf = side_text.text_frame
    stf.word_wrap = True
    items = [
        "Board friendly",
        "Future skills",
        "Student value",
        "Teacher productivity",
        "Admissions story",
    ]
    first = True
    for item in items:
        p = stf.paragraphs[0] if first else stf.add_paragraph()
        first = False
        p.text = item
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PPTX_TEAL
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(16)


if __name__ == "__main__":
    build_pdf()
    build_pptx()
    print(f"Generated {PDF_NAME}")
    print(f"Generated {PPTX_NAME}")
